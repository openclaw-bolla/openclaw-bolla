#!/usr/bin/env python3
"""KI-Folie für Lessing — v3: volle Fläche, große Karten, Hyperlink"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

BG       = RGBColor(0x0D, 0x0B, 0x1E)
ACCENT   = RGBColor(0xFF, 0x66, 0x99)
PURPLE   = RGBColor(0x7C, 0x3A, 0xED)
CYAN     = RGBColor(0x06, 0xB6, 0xD4)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
DIM      = RGBColor(0xB0, 0xA0, 0xD0)
CARD_BG  = RGBColor(0x1E, 0x0F, 0x3A)
CARD_BG2 = RGBColor(0x16, 0x08, 0x2C)
COL_GPT    = RGBColor(0x10, 0xA3, 0x7F)
COL_CLAUDE = RGBColor(0xD4, 0x7A, 0x4A)
COL_GEM    = RGBColor(0x42, 0x85, 0xF4)
COL_META   = RGBColor(0x1B, 0x74, 0xE4)
COL_MS     = RGBColor(0x00, 0xBC, 0xF2)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)
RED        = RGBColor(0xFF, 0x33, 0x44)

def shape(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

def text(slide, l, t, w, h, txt, size, bold=False, italic=False,
         color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tb

def model_card(slide, l, t, w, h, col, name, maker, desc):
    shape(slide, l, t, w, h, CARD_BG2)
    shape(slide, l, t, 50000, h, col)          # linker Streifen
    shape(slide, l, t, w, 24000, col)           # oberer Streifen
    text(slide, l+80000, t+50000,  w-100000, int(h*.28), name, 15, bold=True)
    text(slide, l+80000, t+310000, w-100000, int(h*.18), maker, 11, color=col)
    text(slide, l+80000, t+510000, w-100000, int(h*.38), desc,  11, color=DIM)

def result_card(slide, l, t, w, h, col, icon, title, desc):
    shape(slide, l, t, w, h, CARD_BG)
    shape(slide, l, t, w, 24000, col)
    # Icon-Kreis
    cr = int(h * 0.25); cx = l + (w - cr)//2; cy = t + int(h*.08)
    shape(slide, cx, cy, cr, cr,
          RGBColor(max(0,col[0]-100), max(0,col[1]-100), max(0,col[2]-100)))
    text(slide, cx, cy + int(cr*.05), cr, int(cr*.9), icon, 22, align=PP_ALIGN.CENTER)
    text(slide, l+30000, t+int(h*.42), w-60000, int(h*.28),
         title, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(slide, l+20000, t+int(h*.70), w-40000, int(h*.26),
         desc, 10, color=DIM, align=PP_ALIGN.CENTER)

# ── Folie ─────────────────────────────────────────────────────────────────────
prs = Presentation()
W, H = 12192000, 6858000
prs.slide_width = Emu(W); prs.slide_height = Emu(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])
PAD, GAP = 260000, 55000

shape(slide, 0, 0, W, H, BG)

# ═══ HEADER ═══════════════════════════════════════════════════════════════════
shape(slide, 0, 0, W, 400000, RGBColor(0x12, 0x07, 0x28))
shape(slide, 0, 384000, W, 16000, ACCENT)
text(slide, PAD, 55000, W-2*PAD, 290000,
     "KI-Unterstützung für Lehrerinnen und Lehrer", 24, bold=True)
text(slide, W-3000000, 100000, 2750000, 220000,
     "Lessing Gymnasium · 2026", 11, color=DIM, align=PP_ALIGN.RIGHT)

# ═══ BLOCK 1: ASSISTENT (groß, reicht fast bis zu den Modellen) ═══════════════
y1, b1h = 460000, 1700000
shape(slide, PAD, y1, W-2*PAD, b1h, RGBColor(0x1A, 0x09, 0x35))
shape(slide, PAD, y1, 30000, b1h, ACCENT)

# Label + Erklärungstext links
text(slide, PAD+90000, y1+70000, 4200000, 280000,
     "① Der KI-Assistent", 16, bold=True, color=ACCENT)
text(slide, PAD+90000, y1+340000, 4200000, 240000,
     "z. B. OpenClaw · Claude Code · ChatGPT-App · Gemini Advanced", 13, color=WHITE)
text(slide, PAD+90000, y1+570000, 4200000, 900000,
     "Der Assistent ist das Frontend — die Oberfläche.\n"
     "Das KI-Modell dahinter kann jederzeit gewechselt werden.\n\n"
     "Wie ein Browser: er zeigt Inhalte an — die Website\n(= das Modell) kann eine andere sein.", 12, color=DIM)

# Trennlinie
shape(slide, int(W*.40), y1+80000, 14000, b1h-160000, RGBColor(0x35, 0x18, 0x58))

# 3 Konzept-Boxen rechts — groß, gehen fast bis unten
bx_start = int(W*.41)
bw_total = W - PAD - bx_start - 40000
n_b = 3
bw = int((bw_total - 2*80000) / n_b)
bh = int(b1h * 0.82)
by = y1 + int((b1h - bh) // 2)

boxes = [
    (ACCENT,  "🧑‍💻", "Assistent",  "OpenClaw\nClaude Code\nChatGPT-App"),
    (PURPLE,  "🧠",  "KI-Modell",  "GPT · Claude\nGemini · Llama\nPhi · Mistral"),
    (GREEN,   "✅",  "Ergebnis",   "Text · Plan\nKorrektur\nGrafik · Code"),
]
for i, (col, ic, lbl, sub) in enumerate(boxes):
    bxi = bx_start + i * (bw + 80000)
    shape(slide, bxi, by, bw, bh, col)
    text(slide, bxi, by+int(bh*.06), bw, int(bh*.32), ic, 28, align=PP_ALIGN.CENTER)
    text(slide, bxi+20000, by+int(bh*.42), bw-40000, int(bh*.24),
         lbl, 15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text(slide, bxi+15000, by+int(bh*.64), bw-30000, int(bh*.32),
         sub, 11, color=RGBColor(0xFF,0xFF,0xFF), align=PP_ALIGN.CENTER)
    if i < 2:
        text(slide, bxi+bw+10000, by+int(bh*.38), 60000, int(bh*.24),
             "→", 20, bold=True, color=DIM, align=PP_ALIGN.CENTER)

# ═══ BLOCK 2: MODELLE ═════════════════════════════════════════════════════════
y2 = y1 + b1h + 70000
lbl_h = 230000
text(slide, PAD, y2, W//2, lbl_h,
     "② Die bekanntesten KI-Modelle", 14, bold=True, color=DIM)

y2m = y2 + lbl_h
mod_h = 1550000   # doppelt so hoch wie vorher
models = [
    (COL_GPT,    "GPT-4o",         "OpenAI",    "Text · Code\nVision · Analyse"),
    (COL_GPT,    "GPT-4.5",        "OpenAI",    "Kreativität\nLangtext · Chat"),
    (COL_CLAUDE, "Claude Opus",    "Anthropic", "Reasoning\nLange Dokumente"),
    (COL_CLAUDE, "Claude Sonnet",  "Anthropic", "Schnell &\npräzise"),
    (COL_GEM,    "Gemini 1.5",     "Google",    "Multimodal\nGoogle-Suche"),
    (COL_META,   "Llama 3",        "Meta",      "Open Source\nLokal nutzbar"),
    (COL_MS,     "Phi-3",          "Microsoft", "Klein &\neffizient"),
]
n = len(models)
mod_w = int((W - 2*PAD - (n-1)*GAP) / n)
mx = PAD
for col, name, maker, desc in models:
    model_card(slide, mx, y2m, mod_w, mod_h, col, name, maker, desc)
    mx += mod_w + GAP

# ═══ BLOCK 3: ERGEBNISSE ══════════════════════════════════════════════════════
y3 = y2m + mod_h + 70000
text(slide, PAD, y3, W//2, lbl_h,
     "③ Was kann für Lehrkräfte entstehen?", 14, bold=True, color=DIM)

y3r = y3 + lbl_h
res_h = H - y3r - 160000   # volle verbleibende Höhe nutzen
results = [
    (ACCENT,  "📋", "Unterrichts-\nvorbereitung",  "Stunden & Sequenzen\nin Minuten"),
    (PURPLE,  "✏️",  "Differenzierung",             "Aufgaben auf\n3 Niveaus"),
    (CYAN,    "📊", "Tests &\nPrüfungen",           "Aufgaben + Lösung +\nErwartungshorizont"),
    (GREEN,   "💬", "Feedback",                     "Individuell für\njeden Schüler"),
    (AMBER,   "📧", "Eltern-\nkommunikation",       "Briefe & Mails\nprofessionell"),
    (RGBColor(0xF4,0x72,0x9A), "🌍", "Übersetzungen", "Elterninfos in\n20+ Sprachen"),
    (RGBColor(0x34,0xD3,0x99), "🖼️",  "Bilder &\nGrafiken", "KI-Illustrationen\nfür Materialien"),
    (GOLD,    "⭐", "Korrektur &\nBenotung",        "Foto hochladen →\nNote 1–6 + Feedback"),
]
n_r = len(results)
res_w = int((W - 2*PAD - (n_r-1)*GAP) / n_r)
rx = PAD
for i, (col, icon, title, desc) in enumerate(results):
    result_card(slide, rx, y3r, res_w, res_h, col, icon, title, desc)

    # LIVE DEMO Badge mit Hyperlink auf der letzten Karte
    if i == n_r - 1:
        bw2, bh2 = 370000, 140000
        bx2 = rx + res_w - bw2 - 25000
        by2 = y3r + 30000
        badge = shape(slide, bx2, by2, bw2, bh2, RED)
        tb = text(slide, bx2, by2+18000, bw2, bh2-18000,
                  "▶  LIVE DEMO", 8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Hyperlink auf claude.ai
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            sp = badge._element
            nvSpPr = sp.find(qn('p:nvSpPr'))
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(qn('p:cNvPr'))
                nvPr  = nvSpPr.find(qn('p:nvPr'))
                if nvPr is None:
                    nvPr = etree.SubElement(nvSpPr, qn('p:nvPr'))
                hlinkClick = etree.SubElement(nvPr, qn('a:hlinkClick'))
                rid = slide.part.relate_to(
                    'https://claude.ai', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
                    is_external=True)
                hlinkClick.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', rid)
        except Exception as e:
            pass  # Link klappt nicht → Badge bleibt trotzdem

    rx += res_w + GAP

# Footer
text(slide, PAD, H-155000, W-2*PAD, 140000,
     "💡  KI ist kein Ersatz — sondern ein Werkzeug, das Freiraum für das Wesentliche schafft: den Menschen.",
     10, italic=True, color=DIM, align=PP_ALIGN.CENTER)

out = "/mnt/d/OneDrive/Desktop/KI_fuer_Lehrer.pptx"
prs.save(out)
print(f"Gespeichert: {out}")
