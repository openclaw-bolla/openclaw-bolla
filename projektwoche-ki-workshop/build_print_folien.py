#!/usr/bin/env python3
"""Ergänzt 2 ausdruck-optimierte Folien (helle Versionen von Folie 1 + 2).
   Methode: bestehende Folien duplizieren + Farben per Luminanz/Sättigung
   hell ummappen (dunkler BG -> weiss, weisser Text -> dunkel, Akzente bleiben).
   Originale bleiben unangetastet. Idempotent: baut von 2-Folien-Basis."""
import os, copy, glob, shutil, datetime as _dt
from pptx import Presentation
from pptx.oxml.ns import qn

PPTX = "/mnt/d/OneDrive/Desktop/KI_fuer_Lehrer.pptx"
BACKUP_GLOB = "/home/bolla/workspace/projektwoche-ki-workshop/KI_fuer_Lehrer_backup_*.pptx"

def slide_count(p): return len(Presentation(p).slides)

# Basis mit GENAU 2 Folien (Titel + Meilenstein) bestimmen
if slide_count(PPTX) == 2:
    SRC = PPTX
else:
    SRC = None
    for b in sorted(glob.glob(BACKUP_GLOB), reverse=True):
        try:
            if slide_count(b) == 2: SRC = b; break
        except Exception: pass
    if SRC is None:
        raise SystemExit("Keine 2-Folien-Basis gefunden (Titel+Meilenstein)!")

prs = Presentation(SRC)

# ── Farb-Ummapping dunkel -> hell ──
def remap(r,g,b):
    l = (0.299*r + 0.587*g + 0.114*b)/255          # Luminanz
    mx,mn = max(r,g,b),min(r,g,b)
    sat = (mx-mn)/mx if mx else 0                    # Sättigung
    if l < 0.18:   return (255,255,255)              # dunkler Hintergrund -> weiss
    if sat > 0.40: return (r,g,b)                     # kräftiger Akzent (Pink/Marken) -> belassen
    if l > 0.78:   return (26,26,46)                  # weisser Text -> dunkel
    if l > 0.50:   return (int(r*0.42),int(g*0.40),int(b*0.48))  # heller Subtext -> abdunkeln
    return (r,g,b)

def remap_colors(el):
    for clr in el.iter(qn('a:srgbClr')):
        v = clr.get('val')
        if v and len(v) == 6:
            try: r,g,b = int(v[0:2],16),int(v[2:4],16),int(v[4:6],16)
            except ValueError: continue
            nr,ng,nb = remap(r,g,b)
            clr.set('val', f'{nr:02X}{ng:02X}{nb:02X}')

def dup_slide_print(src_idx):
    source = prs.slides[src_idx]
    blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    new = prs.slides.add_slide(blank)
    # Falls das Blank-Layout selbst Platzhalter/Form mitbringt: entfernen
    for ph in list(new.shapes):
        ph._element.getparent().remove(ph._element)
    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        remap_colors(new_el)
        new.shapes._spTree.append(new_el)
    # Notes übernehmen (falls vorhanden)
    try:
        if source.has_notes_slide and source.notes_slide.notes_text_frame.text.strip():
            new.notes_slide.notes_text_frame.text = "[DRUCK-VERSION] " + source.notes_slide.notes_text_frame.text
    except Exception:
        pass
    return new

# Folie 1 (Titel-Infografik) und Folie 2 (Meilenstein) als Print-Versionen anhängen
dup_slide_print(0)
dup_slide_print(1)

prs.save(PPTX)
print(f"✓ Print-Folien ergänzt. Folien gesamt: {len(prs.slides)} (1-2 Original dunkel, 3-4 Druck hell)")
