#!/usr/bin/env python3
"""Erzeugt die Meilenstein-Folie (Themenblock 1) in KI_fuer_Lehrer.pptx.
   Stil an bestehende Folie angelehnt: dunkel-violett, Segoe UI, Pink-Akzent.
   Idempotent: entfernt eine evtl. vorher erzeugte Meilenstein-Folie (Marker)."""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

PPTX = "/mnt/d/OneDrive/Desktop/KI_fuer_Lehrer.pptx"
BACKUP_GLOB = "/home/bolla/workspace/projektwoche-ki-workshop/KI_fuer_Lehrer_backup_*.pptx"

# ── Palette ──
BG     = RGBColor(0x0D,0x0B,0x1E)
CARD   = RGBColor(0x1E,0x0F,0x3A)
HEADER = RGBColor(0x1A,0x09,0x35)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
SUB    = RGBColor(0xB0,0xA0,0xD0)
PINK   = RGBColor(0xFF,0x66,0x99)
DIMLINE= RGBColor(0x4A,0x3A,0x6A)
FONT   = "Segoe UI"
EMUIN  = 914400

def IN(v): return Emu(int(v*EMUIN))

# Saubere Basis bestimmen (vermeidet verwaiste Slide-Parts / Duplikate):
# - Hat die Zieldatei nur die Titelfolie → sie ist die Basis (+ frisches Backup).
# - Hat sie schon die generierte Folie → vom jüngsten 1-Folien-Backup bauen.
import glob, shutil, datetime as _dt
def slide_count(path):
    return len(Presentation(path).slides)

if os.path.exists(PPTX) and slide_count(PPTX) == 1:
    SRC = PPTX
    shutil.copy2(PPTX, f"/home/bolla/workspace/projektwoche-ki-workshop/KI_fuer_Lehrer_backup_{_dt.datetime.now():%Y%m%d_%H%M}.pptx")
else:
    SRC = None
    for b in sorted(glob.glob(BACKUP_GLOB), reverse=True):
        if slide_count(b) == 1:
            SRC = b; break
    if SRC is None:
        raise SystemExit("Keine saubere 1-Folien-Basis gefunden — bitte Titelfolie sichern!")

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height
MARKER = "MEILENSTEIN_FOLIE_V"

slide = prs.slides.add_slide(prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1])

def rect(l,t,w,h,fill,line=None,round_=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE, IN(l),IN(t),IN(w),IN(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line: shp.line.color.rgb = line; shp.line.width = Pt(1)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def txt(l,t,w,h,text,size,color,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,font=FONT,spacing=None):
    tb = slide.shapes.add_textbox(IN(l),IN(t),IN(w),IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    lines = text.split("\n")
    for idx,ln in enumerate(lines):
        p = tf.paragraphs[0] if idx==0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.line_spacing = spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.name = font; r.font.color.rgb = color
    return tb

# ── Hintergrund ──
rect(0,0,SW/EMUIN,SH/EMUIN,BG)
rect(0,0,SW/EMUIN,0.08,PINK)            # Top-Akzentstreifen

# ── Marker (unsichtbar klein, für Idempotenz) ──
m = txt(0.05,7.42,2,0.1,MARKER+"1",5,BG)

# ── Titel ──
txt(0.5,0.32,12.3,0.5,"④  Vom Buchdruck zur KI", 22, PINK, bold=True)
txt(0.5,0.92,12.3,0.4,"Die Innovationszyklen werden immer kürzer — Technik beschleunigt sich selbst.", 13, SUB)

# ── Zeitstrahl ──
milestones = [
    ("1450","Buchdruck","📜"),
    ("1876","Telefon","☎️"),
    ("1895","Radio","📻"),
    ("1927","Fernsehen","📺"),
    ("1946","Computer","💻"),
    ("1969","Internet","🌐"),
    ("1991","World Wide Web","🔗"),
    ("2007","Smartphone","📱"),
    ("2022","Generative KI","🤖"),
    ("2025","Quanten­computer","⚛️"),
]
n = len(milestones)
left_pad, right_pad = 0.7, 0.7
strip_w = (SW/EMUIN) - left_pad - right_pad
line_y = 2.55
# Zeitstrahl-Linie
ln = rect(left_pad, line_y, strip_w, 0.035, DIMLINE)
step = strip_w / (n-1)
for i,(yr,label,emo) in enumerate(milestones):
    cx = left_pad + i*step
    # Farbverlauf: alt (gedämpft) → neu (pink)
    frac = i/(n-1)
    col = RGBColor(int(0x6A+(0xFF-0x6A)*frac), int(0x5A+(0x66-0x5A)*frac), int(0x8A+(0x99-0x8A)*frac))
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(cx-0.09), IN(line_y-0.075), IN(0.18), IN(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); dot.shadow.inherit=False
    # Emoji über der Linie
    txt(cx-0.5, line_y-0.78, 1.0, 0.4, emo, 19, WHITE, align=PP_ALIGN.CENTER)
    # Jahr + Label unter der Linie (abwechselnd hoch/tief für Lesbarkeit)
    ty = line_y + 0.18
    txt(cx-0.62, ty, 1.24, 0.3, yr, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(cx-0.62, ty+0.28, 1.24, 0.6, label, 9.5, SUB, align=PP_ALIGN.CENTER, spacing=0.9)

# ── Gap-Hinweise (Beschleunigung sichtbar machen) ──
gaps = [
    (0,1,"426 Jahre"),
    (4,5,"23 Jahre"),
    (7,8,"15 Jahre"),
    (8,9,"3 Jahre"),
]
for a,b,lbl in gaps:
    cxa = left_pad + a*step; cxb = left_pad + b*step
    mid = (cxa+cxb)/2
    txt(mid-0.7, line_y-1.12, 1.4, 0.25, "← "+lbl+" →", 8.5, PINK, align=PP_ALIGN.CENTER, bold=True)

# ── Untere Box: Adoptions-Tempo ──
def txt_namejahr(l,t,w,name,jahr):
    tb = slide.shapes.add_textbox(IN(l),IN(t),IN(w),IN(0.3))
    tf = tb.text_frame; tf.word_wrap=False
    tf.margin_left=0; tf.margin_top=0; tf.margin_bottom=0; tf.margin_right=0
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text=name; r1.font.size=Pt(11); r1.font.bold=True; r1.font.name=FONT; r1.font.color.rgb=WHITE
    r2 = p.add_run(); r2.text="  "+jahr; r2.font.size=Pt(8.5); r2.font.name=FONT; r2.font.color.rgb=SUB
    return tb

box_t = 4.25
box_h = 2.85
rect(0.5, box_t, 12.33, box_h, HEADER, round_=True)
txt(0.8, box_t+0.2, 11.7, 0.35, "⚡  Wie schnell erreicht eine Technik 1 Million Nutzer?", 14, WHITE, bold=True)

adopt = [
    ("Netflix","1999","3,5 Jahre", 0.92),
    ("Facebook","2004","10 Monate", 0.55),
    ("Spotify","2008","5 Monate", 0.38),
    ("Instagram","2010","2,5 Monate", 0.22),
    ("ChatGPT","2022","5 Tage", 0.05),
]
bar_left = 2.85
bar_max  = 9.0
row_h = 0.36
y0 = box_t + 0.82
for i,(name,yr,dur,frac) in enumerate(adopt):
    y = y0 + i*row_h
    txt_namejahr(0.8, y+0.01, 1.95, name, yr)
    # Balken (kürzer = schneller); ChatGPT pink hervorgehoben
    w = max(0.25, bar_max*frac)
    col = PINK if name=="ChatGPT" else RGBColor(0x6C,0x5A,0x9C)
    b = rect(bar_left, y+0.02, w, 0.22, col, round_=True)
    txt(bar_left+w+0.14, y, 2.2, 0.26, dur, 10.5, (PINK if name=="ChatGPT" else SUB), bold=(name=="ChatGPT"))

# Fußnote zur Ehrlichkeit der Metrik (unter der letzten Zeile, innerhalb der Box)
txt(0.8, box_t+box_h-0.28, 11.7, 0.25,
    "Hinweis: konsistente Metrik „bis 1 Mio. Nutzer\". Zum Vergleich brauchte das Radio 38 Jahre bis 50 Mio. Hörer.",
    8, RGBColor(0x80,0x72,0xA0))

# ── Speaker Notes (Quellen / Kontext für Chris beim Vortrag) ──
notes = slide.notes_slide.notes_text_frame
notes.text = (
    "Kernbotschaft: Technische Neuerungen verbreiten sich heute in Tagen, früher in Jahrzehnten.\n\n"
    "Zeitstrahl-Daten:\n"
    "1450 Gutenberg-Buchdruck · 1876 Telefon (Bell) · 1895 Radio (Marconi) · 1927 Fernsehen · "
    "1946 ENIAC (erster elektron. Universalrechner) · 1969 ARPANET (Internet-Vorläufer) · "
    "1991 World Wide Web (Tim Berners-Lee) · 2007 iPhone (Smartphone-Ära) · "
    "2022 ChatGPT (generative KI für alle) · ab ca. 2025 Quantencomputer (im Aufkommen).\n\n"
    "Adoptions-Tempo 'bis 1 Mio. Nutzer' (konsistente Metrik): Netflix ~3,5 Jahre, Facebook ~10 Monate, "
    "Spotify ~5 Monate, Instagram ~2,5 Monate, ChatGPT ~5 Tage.\n\n"
    "WICHTIG zur Ehrlichkeit: Die oft zitierten '38 Jahre Radio' beziehen sich auf 50 Mio. Nutzer, "
    "nicht 1 Mio. — daher hier getrennt als Kontext genannt, nicht im direkten Balkenvergleich. "
    "ChatGPT erreichte 100 Mio. Nutzer in ~2 Monaten (schnellste Consumer-App der Geschichte bis dahin)."
)

prs.save(PPTX)
print(f"✓ Folie erstellt. Folien gesamt: {len(prs.slides)}")
